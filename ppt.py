from pptx import Presentation
from pptx.util import Inches
from dotenv import load_dotenv
import os
import google.generativeai as genai
import streamlit as st

# Load environment variables
load_dotenv()

# Configure the generative AI model
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))
model = genai.GenerativeModel("gemini-pro")

def generate_gemini_content(prompt):
    response = model.generate_content([prompt])
    return response.text

def create_presentation(slides_content, file_name="presentation.pptx"):
    prs = Presentation()

    for i, content in enumerate(slides_content):
        slide_layout = prs.slide_layouts[1]  # Use the title and content layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = f"Slide {i + 1}"
        body.text = content

    prs.save(file_name)
    return file_name

# Streamlit interface
st.title("Slide Generator based on Prompt Topic")
prompt = st.text_input("Enter the prompt topic:")

if st.button("Generate Slides"):
    if prompt:
        st.write("Generating content...")
        slide_contents = generate_gemini_content(prompt).split('\n')
        
        st.write("Creating presentation...")
        file_name = create_presentation(slide_contents)
        
        st.write("Presentation created successfully!")
        with open(file_name, "rb") as file:
            btn = st.download_button(
                label="Download Presentation",
                data=file,
                file_name=file_name,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        
        # Display slides
        prs = Presentation(file_name)
        for slide in prs.slides:
            slide_image = slide.shapes.title.text
            st.subheader(slide_image)
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        st.write(run.text)
    else:
        st.error("Please enter a prompt topic.")
